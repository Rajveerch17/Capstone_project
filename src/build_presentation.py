"""Build Bluestock MF capstone presentation (.pptx) from HTML content and dashboard images."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DASHBOARD_DIR = PROJECT_ROOT / "dashboard"
ASSETS_DIR = Path(
    r"C:\Users\rajve\.cursor\projects\c-Users-rajve-Capstone-project\assets"
)
OUTPUT_PATH = PROJECT_ROOT / "Bluestock_MF_Presentation.pptx"

BLUE_DARK = RGBColor(0x0E, 0x4D, 0x92)
BLUE_LIGHT = RGBColor(0x00, 0xA1, 0xDE)
GREEN = RGBColor(0x1E, 0x5B, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLACK = RGBColor(0x33, 0x33, 0x33)

DASHBOARD_PAGES = {
    "page1": ASSETS_DIR
    / "c__Users_rajve_AppData_Roaming_Cursor_User_workspaceStorage_9f1df91e425e0dd44edecb4f1235dc8d_images_image-6eda75eb-c58e-481e-8fd7-9f24c7f37fc1.png",
    "page2": ASSETS_DIR
    / "c__Users_rajve_AppData_Roaming_Cursor_User_workspaceStorage_9f1df91e425e0dd44edecb4f1235dc8d_images_image-4da48821-bcf0-4f85-81f7-b68bedca514b.png",
    "page3": ASSETS_DIR
    / "c__Users_rajve_AppData_Roaming_Cursor_User_workspaceStorage_9f1df91e425e0dd44edecb4f1235dc8d_images_image-3fc3a074-782a-4544-b4ae-d05cfb139614.png",
    "page4": ASSETS_DIR
    / "c__Users_rajve_AppData_Roaming_Cursor_User_workspaceStorage_9f1df91e425e0dd44edecb4f1235dc8d_images_image-f2cc7489-21c7-4fe6-a7ea-6836a2366185.png",
}


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_bullets(slide, items: list[str], left=0.5, top=1.3, width=12.3, height=5.5, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)


def add_section_heading(slide, text: str, top: float) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN


def add_image_fit(slide, image_path: Path, left, top, width, height) -> None:
    if not image_path.exists():
        return
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_kpi_row(slide, kpis: list[tuple[str, str]], top: float) -> None:
    card_w = 2.9
    gap = 0.2
    for i, (value, label) in enumerate(kpis):
        left = 0.5 + i * (card_w + gap)
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(card_w), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        tf = card.text_frame
        tf.text = f"{value}\n{label}"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(11)
            tf.paragraphs[1].font.color.rgb = GRAY
            tf.paragraphs[1].alignment = PP_ALIGN.CENTER


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE_DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        ("BLUESTOCK MUTUAL FUND ANALYTICS", 40, True, WHITE),
        ("End-to-End Data Science Capstone Project", 24, False, RGBColor(0xFF, 0xF3, 0xCD)),
        (
            "Comprehensive Analysis of 40 Fund Schemes | 1M+ Investor Transactions | 4 Years Data",
            18,
            False,
            WHITE,
        ),
        ("", 12, False, WHITE),
        ("Capstone Project Team", 18, True, WHITE),
        ("Complete Pipeline: Day 1-7 Deliverables", 16, False, WHITE),
        ("June 2026", 16, False, RGBColor(0xCC, 0xDD, 0xEE)),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)


def slide_problem_objective(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Problem Statement & Objectives")
    add_section_heading(slide, "The Challenge", 1.2)
    add_bullets(
        slide,
        [
            "With 40+ mutual fund schemes and 1M+ investor transactions, fund houses and investors need data-driven insights for:",
            "Fund selection across multiple risk and return dimensions",
            "Understanding investor behavior patterns and retention risks",
            "Portfolio optimization and sector concentration monitoring",
            "Risk assessment using sophisticated metrics (VaR, CVaR, alpha-beta)",
        ],
        top=1.55,
        height=2.0,
        size=15,
    )
    add_section_heading(slide, "Project Objectives", 3.5)
    add_bullets(
        slide,
        [
            "Build complete ETL pipeline for data ingestion, cleaning, and validation",
            "Compute comprehensive performance metrics (Sharpe, Sortino, alpha-beta)",
            "Develop interactive dashboards for stakeholder visualization",
            "Analyze investor behavior and identify retention risks",
            "Provide data-driven fund recommendations",
            "Create actionable insights for multiple stakeholders",
        ],
        top=3.85,
        height=2.5,
        size=15,
    )
    highlight = slide.shapes.add_shape(1, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xF8)
    highlight.line.color.rgb = BLUE_LIGHT
    tf = highlight.text_frame
    tf.text = "Expected Outcome: Production-ready analytics platform serving fund managers, investors, and regulators."
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK


def slide_data_sources(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Data Sources")
    add_section_heading(slide, "Input Datasets (10 CSV Files)", 1.2)
    add_kpi_row(
        slide,
        [("40", "Fund Schemes"), ("252K+", "NAV Records"), ("1M+", "Transactions"), ("2K+", "Holdings")],
        1.65,
    )
    add_section_heading(slide, "Data Pipeline Stages", 3.0)
    stages = ["Day 1\nIngestion", "Day 2\nCleaning", "Day 3\nEDA", "Day 4\nPerformance", "Day 5\nDashboard", "Day 6\nAnalytics"]
    for i, stage in enumerate(stages):
        left = 0.5 + i * 2.05
        box = slide.shapes.add_shape(1, Inches(left), Inches(3.35), Inches(1.85), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xF8)
        box.line.color.rgb = BLUE_LIGHT
        tf = box.text_frame
        tf.text = stage
        for p in tf.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK
            p.alignment = PP_ALIGN.CENTER
    add_section_heading(slide, "Technology Stack", 4.5)
    add_bullets(
        slide,
        [
            "Languages: Python 3.8+  |  Databases: SQLite",
            "Data Processing: Pandas, NumPy, SciPy",
            "Visualization: Plotly, Matplotlib, Seaborn",
            "Notebooks: Jupyter Lab",
        ],
        top=4.85,
        height=2.0,
        size=15,
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Architecture")
    add_section_heading(slide, "End-to-End Pipeline", 1.2)
    add_bullets(
        slide,
        [
            "Master Orchestrator: run_pipeline.py executes all 6 stages or runs individual stages",
            "python run_pipeline.py  →  Run full pipeline (Day 1-6)",
            "python run_pipeline.py --stage day4  →  Run specific stage",
            "python run_pipeline.py --list  →  List all stages",
        ],
        top=1.55,
        height=1.8,
        size=14,
    )
    add_section_heading(slide, "Directory Structure", 3.3)
    add_bullets(
        slide,
        [
            "data/raw/  →  Original 10 CSV datasets",
            "data/processed/  →  Cleaned datasets (output)",
            "src/  →  Python pipeline modules + orchestrator",
            "dashboard/  →  Interactive HTML dashboards",
            "notebooks/  →  Jupyter analysis notebooks",
            "reports/  →  Daily summary reports",
            "sql/  →  Database schema & queries",
        ],
        top=3.65,
        height=2.2,
        size=13,
    )
    add_section_heading(slide, "Output Artifacts", 5.7)
    add_bullets(
        slide,
        [
            "CSV Reports: fund_scorecard.csv, var_cvar_report.csv, cohort_analysis.csv",
            "Database: bluestock_mf.db (SQLite with normalized schema)",
            "Dashboards: 4-page interactive dashboard with 15+ visualizations",
            "Documentation: Comprehensive README + final report",
        ],
        top=6.0,
        height=1.2,
        size=13,
    )


def slide_eda_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "EDA Highlights — Market Trends & Performance")
    add_bullets(
        slide,
        [
            "2023 Bull Run: All 40 funds appreciated 20-35% with strong inflows",
            "2024 Correction: Market correction resulted in 8-12% peak-to-trough declines",
            "2025 Stabilization: Recovery phase with stable growth trajectory",
        ],
        top=1.2,
        width=5.8,
        height=1.5,
        size=13,
    )
    add_kpi_row(
        slide,
        [("₹2M+ Cr", "Total Industry AUM"), ("14.5%", "Average CAGR"), ("8", "Top AMC Schemes"), ("₹1.2K Cr", "Monthly SIP Inflows")],
        2.75,
    )
    add_image_fit(slide, DASHBOARD_DIR / "page1_aum_trend.png", 0.4, 3.95, 6.2, 3.2)
    add_image_fit(slide, DASHBOARD_DIR / "page4_sip_nifty_trend.png", 6.8, 3.95, 6.0, 3.2)


def slide_eda_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "EDA Highlights — Investor Demographics")
    add_kpi_row(
        slide,
        [("65%", "SIP Dominance"), ("40%+", "Metro Concentration"), ("25-45", "Peak Age Group"), ("98%", "KYC Compliant")],
        1.2,
    )
    add_bullets(
        slide,
        [
            "Tier-1 Cities: 60% of inflows (Delhi, Mumbai, Bangalore dominate)",
            "Tier-2 Cities: 25% growth opportunity in emerging metros",
            "Tier-3 & Beyond: 15% untapped potential for expansion",
            "Category Mix: Large Cap (18%) | Mid Cap (22%) | Small Cap (8%) | Balanced (28%) | Debt (24%)",
        ],
        top=2.4,
        width=12.3,
        height=1.5,
        size=13,
    )
    add_image_fit(slide, DASHBOARD_DIR / "page3_transaction_type_split.png", 0.4, 3.9, 4.0, 3.3)
    add_image_fit(slide, DASHBOARD_DIR / "page3_amount_by_state.png", 4.6, 3.9, 4.0, 3.3)
    add_image_fit(slide, DASHBOARD_DIR / "page3_avg_sip_by_age_group.png", 8.8, 3.9, 4.0, 3.3)


def slide_performance_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Performance Metrics — Fund Rankings")
    add_bullets(
        slide,
        [
            "Returns: 1yr, 3yr, 5yr CAGR  |  Risk-Adjusted: Sharpe, Sortino",
            "Market Sensitivity: Alpha & Beta vs NIFTY100  |  Downside Risk: Max Drawdown, VaR, CVaR",
            "1. HDFC Top 100 Direct — Score: 82.5 | 3yr CAGR: 19.2% | Sharpe: 1.62",
            "2. SBI Bluechip — Score: 81.2 | 3yr CAGR: 18.8% | Sharpe: 1.58",
            "3. ICICI Bluechip — Score: 79.8 | 3yr CAGR: 18.5% | Sharpe: 1.55",
        ],
        top=1.2,
        width=12.3,
        height=2.0,
        size=13,
    )
    add_image_fit(slide, DASHBOARD_DIR / "page2_return_vs_risk.png", 0.4, 3.3, 6.2, 3.8)
    add_image_fit(slide, DASHBOARD_DIR / "page2_nav_vs_benchmark.png", 6.8, 3.3, 6.0, 3.8)


def slide_performance_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Performance Metrics — Risk Analysis & Scorecard")
    weights = [("30%", "3yr CAGR"), ("25%", "Sharpe Ratio"), ("20%", "Alpha"), ("15%", "Expense Ratio"), ("10%", "Max Drawdown")]
    for i, (pct, label) in enumerate(weights):
        left = 0.4 + i * 2.5
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.25), Inches(2.3), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        tf = box.text_frame
        tf.text = f"{pct}\n{label}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(10)
            tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    add_kpi_row(
        slide,
        [("12", "High Sharpe (>1.5)"), ("28", "Positive Alpha"), ("8", "Low Vol (<12%)"), ("70%", "Beat NIFTY100")],
        2.3,
    )
    add_image_fit(slide, DASHBOARD_DIR / "page2_scorecard_table.png", 0.4, 3.55, 12.3, 3.6)


def slide_dashboard_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Dashboard Screenshots — Page 1 & 2")
    add_bullets(
        slide,
        [
            "Page 1: Industry Overview — KPIs, AUM trend, AMC rankings",
            "Page 2: Fund Performance — Risk-return scatter, scorecard table, NAV vs benchmark",
        ],
        top=1.15,
        height=0.8,
        size=13,
    )
    add_image_fit(slide, DASHBOARD_PAGES["page1"], 0.35, 1.85, 6.2, 5.3)
    add_image_fit(slide, DASHBOARD_PAGES["page2"], 6.75, 1.85, 6.2, 5.3)


def slide_dashboard_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Dashboard Screenshots — Page 3 & 4")
    add_bullets(
        slide,
        [
            "Page 3: Investor Demographics — Geographic heat map, age groups, transaction types",
            "Page 4: Category & Sector Trends — SIP inflows, NIFTY50 trend, category heatmap",
        ],
        top=1.15,
        height=0.8,
        size=13,
    )
    add_image_fit(slide, DASHBOARD_PAGES["page3"], 0.35, 1.85, 6.2, 5.3)
    add_image_fit(slide, DASHBOARD_PAGES["page4"], 6.75, 1.85, 6.2, 5.3)


def slide_key_findings(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Findings & Strategic Recommendations")
    add_section_heading(slide, "Critical Insights", 1.2)
    add_bullets(
        slide,
        [
            "Fund Selection: Top 12 funds show Sharpe ratios >1.5; composite scoring outperforms single-metric ranking",
            "Risk Profile: Average VaR 95% = -2.8% daily; CVaR (tail risk) = -4.1%; max drawdown 8-22%",
            "Investor Behavior: SIP dominance (65%); metro concentration (40%); peak age 25-45; excellent KYC compliance",
            "SIP Continuity Risk: 8% of SIPs show >35-day gaps; early cohorts (pre-2022) show 85% continuity",
            "Sector Concentration: Moderate HHI (avg 0.18); IT/Finance/FMCG dominant",
        ],
        top=1.55,
        height=2.8,
        size=13,
    )
    add_section_heading(slide, "Strategic Recommendations", 4.3)
    add_bullets(
        slide,
        [
            "For Fund Managers: Engage at-risk SIPs proactively; develop Tier-2/3 distribution; optimize costs",
            "For Investors: Use composite scorecard; diversify across categories/sectors; maintain SIP discipline",
            "For Regulators: Monitor concentration risk (HHI); enhance drawdown disclosures; standardize reporting",
        ],
        top=4.65,
        height=2.0,
        size=13,
    )


def slide_thank_you(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE_DARK)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    content = [
        ("Thank You!", 44, True, WHITE),
        ("Questions & Discussion", 22, False, RGBColor(0xFF, 0xF3, 0xCD)),
        ("", 10, False, WHITE),
        ("Project Summary", 18, True, WHITE),
        ("Complete 7-day capstone execution (Day 1-7 deliverables)", 14, False, WHITE),
        ("ETL pipeline + Analytics + Interactive 4-page dashboard", 14, False, WHITE),
        ("40 fund schemes, 1M+ transactions, 252K+ NAV records analyzed", 14, False, WHITE),
        ("Production-ready codebase with comprehensive documentation", 14, False, WHITE),
        ("", 10, False, WHITE),
        ("Deliverables: GitHub Repository | Final Report | README Setup Guide", 14, False, RGBColor(0xCC, 0xDD, 0xEE)),
        ("Version 1.0 | June 2026", 12, False, RGBColor(0xAA, 0xBB, 0xCC)),
    ]
    for i, (text, size, bold, color) in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem_objective(prs)
    slide_data_sources(prs)
    slide_architecture(prs)
    slide_eda_1(prs)
    slide_eda_2(prs)
    slide_performance_1(prs)
    slide_performance_2(prs)
    slide_dashboard_1(prs)
    slide_dashboard_2(prs)
    slide_key_findings(prs)
    slide_thank_you(prs)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(f"Created presentation: {path}")
